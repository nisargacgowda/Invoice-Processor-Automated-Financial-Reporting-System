from pptx import Presentation
from pptx.util import Inches
import pandas as pd

def generate_ppt():

    df = pd.read_excel(
        "all_invoices.xlsx"
    )

    total_spending = df["Amount"].sum()

    ppt = Presentation()

    slide = ppt.slides.add_slide(
        ppt.slide_layouts[1]
    )

    slide.shapes.title.text = \
        "Invoice Analysis Report"

    slide.placeholders[1].text = \
        f"Total Spending: ₹{total_spending}"

    slide.shapes.add_picture(
        "vendor_chart.png",
        Inches(1),
        Inches(2),
        width=Inches(5)
    )

    ppt.save(
        "Invoice_Report.pptx"
    )

    print("PPT Generated")

if __name__ == "__main__":
    generate_ppt()